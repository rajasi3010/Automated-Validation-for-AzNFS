#!/usr/bin/env python3
"""Generate the AzNFS Automated Validation demo deck (~10 min, team audience)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette (Azure) ------------------------------------------------------
AZURE   = RGBColor(0x00, 0x78, 0xD4)
DARK    = RGBColor(0x1B, 0x1B, 0x2F)
SLATE   = RGBColor(0x3B, 0x3B, 0x4F)
LIGHT   = RGBColor(0xF3, 0xF6, 0xFB)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN   = RGBColor(0x10, 0x7C, 0x10)
AMBER   = RGBColor(0xC1, 0x7A, 0x00)
GREY    = RGBColor(0x60, 0x60, 0x70)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb


def set_text(tf, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for item in runs:
        text, size, color, bold = item[0], item[1], item[2], item[3]
        bullet = item[4] if len(item) > 4 else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space)
        p.space_before = Pt(0)
        r = p.add_run(); r.text = ("•  " + text) if bullet else text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Segoe UI"


def bg(s, color):
    r = rect(s, 0, 0, SW, SH, color)
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)


def header(s, title, kicker=None):
    bg(s, WHITE)
    rect(s, 0, 0, SW, Inches(1.15), DARK)
    rect(s, 0, Inches(1.15), SW, Pt(4), AZURE)
    if kicker:
        kb = box(s, Inches(0.6), Inches(0.18), Inches(12), Inches(0.35))
        set_text(kb.text_frame, [(kicker, 12, AZURE, True)])
    tb = box(s, Inches(0.6), Inches(0.42), Inches(12.1), Inches(0.65))
    set_text(tb.text_frame, [(title, 28, WHITE, True)])


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ===========================================================================
# 1 — TITLE
# ===========================================================================
s = slide(); bg(s, DARK)
rect(s, 0, Inches(3.05), SW, Pt(4), AZURE)
t = box(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.1))
set_text(t.text_frame, [("Automated Validation for AzNFS", 44, WHITE, True)])
st = box(s, Inches(0.9), Inches(3.25), Inches(11.5), Inches(0.8))
set_text(st.text_frame, [("Zero-touch distro support for Azure Files NFS", 22, AZURE, False)])
who = box(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.9))
set_text(who.text_frame, [
    ("Rajasi Mandal  ·  Azure Files", 16, WHITE, False),
    ("Team demo  ·  bi-weekly sync", 13, GREY, False),
])
notes(s, "One line: today, adding AzNFS support for a new Linux distro is a manual chore. "
         "This tool makes it fully automatic — discover, check, validate on a real VM — end to end. "
         "Keep it to ~10 min, high level.")

# ===========================================================================
# 2 — THE PROBLEM
# ===========================================================================
s = slide(); header(s, "The problem: supporting a new distro is manual", "WHY WE BUILT THIS")
tb = box(s, Inches(0.7), Inches(1.6), Inches(7.2), Inches(5.3))
set_text(tb.text_frame, [
    ("New Linux distros & versions ship all the time", 20, DARK, True, True),
    ("Each one needs AzNFS support added by hand:", 16, SLATE, False),
    ("build → publish the package → manually test the mount", 15, GREY, False),
    ("We're busy with other work, so support lags behind", 20, DARK, True, True),
    ("Customers end up asking us:", 16, SLATE, False),
    ("\u201cWhy isn\u2019t my distro supported for this package?\u201d", 17, AMBER, True),
    ("No single, always-current view of what\u2019s validated", 20, DARK, True, True),
])
# callout card
rect(s, Inches(8.2), Inches(1.9), Inches(4.4), Inches(4.3), LIGHT)
rect(s, Inches(8.2), Inches(1.9), Inches(0.12), Inches(4.3), AMBER)
c = box(s, Inches(8.5), Inches(2.2), Inches(3.9), Inches(3.8))
set_text(c.text_frame, [
    ("Today", 15, AMBER, True),
    ("Manual · reactive · slow", 18, DARK, True),
    ("", 6, DARK, False),
    ("Effort scales with every new distro release, and gaps are only found when a customer hits them.", 15, SLATE, False),
])
notes(s, "Frame the pain: manual effort, we're busy, releases lag, and customers surface the gaps for us. "
         "This is reactive and doesn't scale as distros multiply.")

# ===========================================================================
# 3 — THE IDEA
# ===========================================================================
s = slide(); header(s, "The idea: make it a hands-off pipeline", "WHAT THIS TOOL DOES")
big = box(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.7))
set_text(big.text_frame, [
    ("Automatically discover every new distro release, check the AzNFS "
     "package, and validate the mount on a real Azure VM \u2014 on a schedule, "
     "with email alerts.", 24, DARK, True),
])
# three value chips
chips = [
    ("Proactive", "Finds new distros before customers ask", GREEN),
    ("Hands-off", "Humans step in only when something fails", AZURE),
    ("Trustworthy", "Every verdict comes from a real VM test", SLATE),
]
x = Inches(0.9)
for title, sub, col in chips:
    rect(s, x, Inches(4.3), Inches(3.7), Inches(2.1), LIGHT)
    rect(s, x, Inches(4.3), Inches(3.7), Inches(0.5), col)
    hb = box(s, x + Inches(0.15), Inches(4.35), Inches(3.4), Inches(0.5))
    set_text(hb.text_frame, [(title, 16, WHITE, True)])
    bb = box(s, x + Inches(0.2), Inches(5.0), Inches(3.3), Inches(1.3))
    set_text(bb.text_frame, [(sub, 15, SLATE, False)])
    x += Inches(3.95)
notes(s, "The one-liner. Scope is Azure Files NFS only. Three payoffs: proactive, hands-off, trustworthy.")

# ===========================================================================
# 4 — HOW IT WORKS (architecture)
# ===========================================================================
s = slide(); header(s, "How it works: a 3-phase pipeline", "ARCHITECTURE")
stages = [
    ("Azure\nMarketplace", GREY, "source"),
    ("Phase 1\nDiscover", AZURE, "Scan for new\ndistro releases"),
    ("Phase 2\nGate & Check", AZURE, "Is AzNFS published\n& newer?"),
    ("Phase 3\nValidate", AZURE, "Deploy real VM,\nmount & test"),
    ("Verdict\n+ Email", GREEN, "supported /\nunsupported"),
]
x = Inches(0.55); y = Inches(2.3); w = Inches(2.15); h = Inches(1.5)
gap = Inches(0.35)
for i, (title, col, sub) in enumerate(stages):
    rect(s, x, y, w, h, col)
    tf = box(s, x, y + Inches(0.2), w, Inches(1.1))
    set_text(tf.text_frame, [(title, 15, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    sb = box(s, x, y + h + Inches(0.1), w, Inches(1.0))
    set_text(sb.text_frame, [(sub, 12, GREY, False)], align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        ar = box(s, x + w, y + Inches(0.5), gap, Inches(0.6))
        set_text(ar.text_frame, [("\u2192", 24, AZURE, True)], align=PP_ALIGN.CENTER)
    x += w + gap
# footer band
rect(s, Inches(0.55), Inches(5.4), Inches(12.2), Inches(1.3), LIGHT)
fb = box(s, Inches(0.85), Inches(5.55), Inches(11.6), Inches(1.05))
set_text(fb.text_frame, [
    ("Runs unattended on GitHub Actions (self-hosted Azure VM). Phases chain automatically "
     "and share one state database, so the pipeline always knows what\u2019s been validated.", 15, SLATE, False),
    ("Safety nets: monthly status digest  ·  watchdog alerts if a scheduled run is missed", 14, GREY, True),
])
notes(s, "Walk left to right once. Don't deep dive — this is the mental model. "
         "Chained on GitHub Actions, shared DB is the source of truth.")

# ===========================================================================
# helper for phase slides
# ===========================================================================
def phase_slide(num, name, tag, points, outcome):
    s = slide(); header(s, f"Phase {num}: {name}", tag)
    rect(s, Inches(0.7), Inches(1.55), Inches(1.5), Inches(1.5), AZURE)
    nb = box(s, Inches(0.7), Inches(1.55), Inches(1.5), Inches(1.5))
    set_text(nb.text_frame, [(str(num), 54, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb = box(s, Inches(2.5), Inches(1.55), Inches(10), Inches(3.7))
    runs = [(p, 18, DARK, False, True) for p in points]
    set_text(tb.text_frame, runs, space=12)
    rect(s, Inches(0.7), Inches(5.6), Inches(12), Inches(1.1), LIGHT)
    rect(s, Inches(0.7), Inches(5.6), Inches(0.12), Inches(1.1), GREEN)
    ob = box(s, Inches(1.0), Inches(5.75), Inches(11.5), Inches(0.85))
    set_text(ob.text_frame, [("Result:  ", 15, GREEN, True), (outcome, 15, SLATE, False)]
             if False else [("Result", 13, GREEN, True), (outcome, 16, SLATE, False)])
    return s


# 5 — Phase 1
s = phase_slide(1, "Discover new distros", "SCAN",
    ["Scans Azure Marketplace on a daily schedule",
     "Tracks every distro image in a state DB (new / updated / unchanged)",
     "Emails the team only when a genuinely new distro release appears"],
    "New releases are captured and handed to Phase 2 \u2014 no manual watching of the Marketplace.")
notes(s, "Phase 1 = the eyes. Daily scan, remembers what it has seen, alerts only on truly new releases.")

# 6 — Phase 2
s = phase_slide(2, "Gate & publish-check", "DECIDE",
    ["For each new distro, checks the AzNFS package on packages.microsoft.com",
     "Version-aware: only proceeds if AzNFS is newer than what we last validated",
     "Decides per distro: needs validation · already good · not supported"],
    "Only the distros that actually need a fresh mount test move on to Phase 3.")
notes(s, "Phase 2 = the gatekeeper. Confirms the package exists and is new enough, avoids wasted VM runs.")

# 7 — Phase 3
s = phase_slide(3, "Validate on a real VM", "PROVE",
    ["LISA spins up an actual Azure VM of that distro",
     "Installs AzNFS, mounts an Azure Files NFS share, runs functional + resilience tests",
     "Auto-records the verdict and emails a summary of what passed / failed"],
    "Distro is marked supported or unsupported automatically \u2014 a real, repeatable proof, not a guess.")
notes(s, "Phase 3 = the proof. Real VM, real mount, real tests via LISA. Verdict is written back automatically.")

# ===========================================================================
# 8 — IMPACT
# ===========================================================================
s = slide(); header(s, "Why it matters", "IMPACT")
items = [
    ("Faster support", "New distros validated automatically, no manual toil", GREEN),
    ("Proactive", "We close gaps before customers hit them", AZURE),
    ("Always current", "One trustworthy list of validated distros", SLATE),
    ("Resilient", "Monthly digest + watchdog for missed runs", AMBER),
]
x = Inches(0.7); y = Inches(1.8)
for i, (title, sub, col) in enumerate(items):
    cx = x + (Inches(6.15) * (i % 2))
    cy = y + (Inches(2.35) * (i // 2))
    rect(s, cx, cy, Inches(5.85), Inches(2.05), LIGHT)
    rect(s, cx, cy, Inches(0.5), Inches(2.05), col)
    hb = box(s, cx + Inches(0.75), cy + Inches(0.25), Inches(4.9), Inches(0.6))
    set_text(hb.text_frame, [(title, 20, DARK, True)])
    bb = box(s, cx + Inches(0.75), cy + Inches(0.95), Inches(4.9), Inches(1.0))
    set_text(bb.text_frame, [(sub, 15, SLATE, False)])
notes(s, "The payoff for the team: faster, proactive, single source of truth, self-monitoring.")

# ===========================================================================
# 9 — WHAT'S NEXT / CLOSE
# ===========================================================================
s = slide(); bg(s, DARK)
rect(s, 0, Inches(1.15), SW, Pt(4), AZURE)
tb = box(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.9))
set_text(tb.text_frame, [("What\u2019s next", 30, WHITE, True)])
pts = box(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(3.2))
set_text(pts.text_frame, [
    ("Productize it and roll out for the team", 22, WHITE, False, True),
    ("Broaden distro coverage as new releases land", 22, WHITE, False, True),
    ("Fold into our regular release process so support never lags", 22, WHITE, False, True),
], space=16)
rect(s, Inches(0.9), Inches(5.6), Inches(11.5), Pt(3), AZURE)
tk = box(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.9))
set_text(tk.text_frame, [("Thank you  ·  Questions?", 22, AZURE, True)])
notes(s, "Close: let's productize this soon so distro support is never a manual bottleneck again. Open for questions.")

out = "AzNFS-Automated-Validation-Demo.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
